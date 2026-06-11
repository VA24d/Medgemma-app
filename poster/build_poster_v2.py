"""
Med Veda — SparQ 2026 Research Poster v2
Fixed: text overflow, header/box clipping, proper images.
"""

import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(REPO, "Sparq-2026-Poster Format.pptx")
OUT_PPTX = os.path.join(REPO, "poster", "Med-Veda-SparQ-2026-Poster-v2.pptx")
ASSETS = os.path.join(REPO, "poster", "assets")

# Colors
DARK_NAVY    = RGBColor(0x0A, 0x1A, 0x2F)
QC_BLUE      = RGBColor(0x00, 0x5A, 0x9C)
TIER1_GREEN  = RGBColor(0x1B, 0x5E, 0x20)
TIER2_BLUE   = RGBColor(0x15, 0x65, 0xC0)
TIER3_PURPLE = RGBColor(0x45, 0x27, 0xA0)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x00, 0x00, 0x00)
ACCENT_GOLD  = RGBColor(0xFF, 0xB3, 0x00)
SECTION_BG   = RGBColor(0xF0, 0xF4, 0xF8)
BORDER_BLUE  = RGBColor(0x00, 0x7A, 0xCC)
TRUST_BG     = RGBColor(0xE8, 0xF5, 0xE9)
KPI_BG       = RGBColor(0xFD, 0xF2, 0xE9)
HIGHLIGHT_BG = RGBColor(0xE3, 0xF2, 0xFD)
ORANGE       = RGBColor(0xE6, 0x51, 0x00)
LIGHT_GRAY   = RGBColor(0xF5, 0xF5, 0xF5)


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_w=Pt(2)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    return shape


def _set_run(run, text, size, bold, color, name="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_textbox(slide, left, top, width, height, text, font_size=22,
                bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    _set_run(p.add_run(), text, font_size, bold, color, font_name)
    return txBox


def section_block(slide, left, top, width, height, title, title_bg, body_bg, border_color):
    """Draw a section: header bar (flat-bottom rectangle) sitting flush on top of the body box."""
    hdr_h = 0.5
    body_top = top + hdr_h

    # Body box (full height minus header — flat top corners hidden behind header)
    add_rect(slide, left, body_top - 0.05, width, height - hdr_h + 0.05,
             body_bg, border_color, Pt(2))

    # Header bar overlaps body top edge
    add_rect(slide, left, top, width, hdr_h, title_bg, border_color, Pt(2))

    # Title text (well inside the bar)
    add_textbox(slide, left + 0.2, top + 0.05, width - 0.4, hdr_h - 0.1,
                title, font_size=28, bold=True, color=WHITE)

    return body_top + 0.15  # content start y


def add_para(tf, text, size=18, bold=False, color=BLACK, space_after=4, name="Calibri"):
    p = tf.add_paragraph()
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    _set_run(p.add_run(), text, size, bold, color, name)
    return p


def build_poster():
    # Extract logo
    tmpl = Presentation(TEMPLATE)
    logo_path = None
    for shape in tmpl.slides[0].shapes:
        if shape.name == "Picture 2":
            logo_path = os.path.join(ASSETS, "sparq_logo.png")
            with open(logo_path, "wb") as f:
                f.write(shape.image.blob)
            break

    # New presentation 24x36
    prs = Presentation()
    prs.slide_width = Inches(24)
    prs.slide_height = Inches(36)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    add_rect(slide, 0, 0, 24, 36, WHITE)

    # ================================================================
    # TOP BANNER
    # ================================================================
    if logo_path:
        slide.shapes.add_picture(logo_path, Inches(0), Inches(0),
                                 Inches(8.4), Inches(5.36))

    add_rect(slide, 7.64, 0, 16.36, 4.2, DARK_NAVY)

    add_textbox(slide, 8.2, 0.35, 15.0, 1.5,
                "Med Veda", font_size=84, bold=True, color=WHITE)
    add_textbox(slide, 8.2, 1.8, 15.0, 0.8,
                "Pervasive Clinical AI on Snapdragon + Governed Hybrid Compute",
                font_size=34, bold=False, color=ACCENT_GOLD)
    add_textbox(slide, 8.2, 2.6, 15.0, 0.9,
                "On-device MedGemma 1.5 4B multimodal  |  Edge companion hub  |  "
                "Optional cloud APIs  |  Privacy-first AI as a continuum",
                font_size=20, color=RGBColor(0xBB, 0xBB, 0xBB))

    # Status chips
    chip_y = 3.55
    for cx, label, clr in [(8.2, "SHIPPED", TIER1_GREEN),
                           (10.7, "PROTOTYPE", TIER2_BLUE),
                           (13.2, "ROADMAP", TIER3_PURPLE)]:
        add_rounded_rect(slide, cx, chip_y, 2.2, 0.38, clr)
        add_textbox(slide, cx, chip_y + 0.04, 2.2, 0.3,
                    label, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Divider
    add_rect(slide, 0, 4.85, 24, 0.08, QC_BLUE)

    # ================================================================
    # COLUMN LAYOUT constants
    # ================================================================
    L_LEFT = 0.5          # left column x
    L_W = 7.8             # left column width
    R_LEFT = 8.8          # right column x
    R_W = 14.7            # right column width
    GAP = 0.35            # vertical gap between sections
    PAD = 0.25            # inner padding

    # ================================================================
    # LEFT COLUMN
    # ================================================================
    # --- Abstract ---
    y = 5.15
    abs_h = 6.6
    cy = section_block(slide, L_LEFT, y, L_W, abs_h,
                        "Abstract", QC_BLUE, SECTION_BG, BORDER_BLUE)

    abs_text = (
        "Mobile clinical workflows operate under intermittent connectivity, strict "
        "latency budgets at the point of care, and thermal constraints on sustained "
        "inference. Protected health information must remain amenable to privacy-"
        "preserving processing, yet multimodal decisions still require tight coupling "
        "between longitudinal records and high-resolution imaging.\n\n"
        "Med Veda is a deployable Android clinician assistant centered on on-device "
        "Google MedGemma 1.5 4B multimodal inference. The system architecture "
        "partitions responsibility explicitly: privacy-sensitive longitudinal reasoning, "
        "chart-conditioned dialogue, and multilingual clinical generation execute "
        "locally on Snapdragon-class hardware behind a native GGUF runtime suitable "
        "for heterogeneous acceleration including NPU-assisted paths. This edge-first "
        "design preserves low-latency interaction in offline-first wards and keeps "
        "narrative PHI under local control.\n\n"
        "Compute-intensive vision stages are selectively orchestrated to enterprise "
        "imaging services under institutional credentials, returning compact structured "
        "analytics rather than open-ended cloud language generation. FHIR-oriented "
        "export enables hospital on-premises or VPC GPU clusters for deeper screening "
        "and cohort analytics while the handset remains the bedside copilot.\n\n"
        "The contribution is a production-oriented hybrid edge-cloud pattern for "
        "pervasive clinical AI: Snapdragon-edge multimodal reasoning as the default "
        "trust boundary, governed delegation for heavy vision, and standards-based "
        "extensibility for institutional compute."
    )
    add_textbox(slide, L_LEFT + PAD, cy, L_W - PAD * 2, abs_h - 0.85,
                abs_text, font_size=16, color=BLACK)

    # --- Problem ---
    y += abs_h + GAP
    prob_h = 4.6
    cy = section_block(slide, L_LEFT, y, L_W, prob_h,
                        "Problem Statement", RGBColor(0xC6, 0x28, 0x28), SECTION_BG, RGBColor(0xC6, 0x28, 0x28))

    problems = [
        "Clinicians need longitudinal chart + imaging at the bedside, often with intermittent or no connectivity",
        "Cloud-bound AI introduces unacceptable latency at point of care and mandates PHI egress to third-party servers",
        "India's DPDP Act and global data protection laws impose strict constraints on clinical data movement",
        "A single device cannot handle all workloads at full speed — large imaging studies hit thermal and compute limits",
        "Existing solutions force a false choice: \"all on phone\" (slow) or \"all in cloud\" (privacy-violating, offline-incompatible)",
    ]
    txBox = slide.shapes.add_textbox(
        Inches(L_LEFT + PAD), Inches(cy), Inches(L_W - PAD * 2), Inches(prob_h - 0.85))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.paragraphs[0].text = ""
    for item in problems:
        add_para(tf, f"•  {item}", size=16, color=BLACK, space_after=6)

    # --- Methods ---
    y += prob_h + GAP
    meth_h = 6.5
    cy = section_block(slide, L_LEFT, y, L_W, meth_h,
                        "Methods & Technical Approach", QC_BLUE, SECTION_BG, BORDER_BLUE)

    methods = [
        ("MedGemma 1.5 4B IT (Multimodal)", "Deployed via GGUF (Q4_K_M default, 2.4 GB) through custom llama.cpp JNI wrapper (aichatlib) with lazy mmproj vision encoder loading"),
        ("Heterogeneous Acceleration", "CPU (KleidiAI on arm64), Adreno GPU, and Hexagon NPU runtime paths; user-selectable energy mode with 2B/4B dynamic switching"),
        ("Chart-Grounded Inference", "Room DB stores longitudinal patient records; system prompts inject full chart context with dated entries for temporally-aware reasoning"),
        ("Hybrid Vision Pipeline", "Study pixels sent to near-edge GPU (prototype: RTX 4070); only structured JSON findings return to device for local text synthesis"),
        ("FHIR Interoperability", "On-device LLM-assisted FHIR bundle export enables hospital VPC ingestion for batch analytics with larger models"),
    ]
    txBox = slide.shapes.add_textbox(
        Inches(L_LEFT + PAD), Inches(cy), Inches(L_W - PAD * 2), Inches(meth_h - 0.85))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.paragraphs[0].text = ""
    for title, desc in methods:
        p = tf.add_paragraph()
        p.space_after = Pt(2)
        p.space_before = Pt(6)
        r = p.add_run()
        _set_run(r, title, 17, True, QC_BLUE)
        p2 = tf.add_paragraph()
        p2.space_after = Pt(6)
        r2 = p2.add_run()
        _set_run(r2, desc, 15, False, BLACK)

    # --- Three-Tier Table ---
    y += meth_h + GAP
    tier_h = 6.2
    cy = section_block(slide, L_LEFT, y, L_W, tier_h,
                        "Three-Tier Pervasive AI Architecture", DARK_NAVY, SECTION_BG, DARK_NAVY)

    tiers = [
        ("Tier 1 — Device Edge", "SHIPPED", TIER1_GREEN, RGBColor(0xE8, 0xF5, 0xE9),
         "Android phone (Snapdragon): MedGemma 1.5 4B GGUF + mmproj; chart-grounded chat; Room DB; PIN auth; full offline operation"),
        ("Tier 2 — Near Edge", "PROTOTYPE", TIER2_BLUE, RGBColor(0xE3, 0xF2, 0xFD),
         "Hospital / clinic LAN GPU (lab: RTX 4070): Heavy imaging → structured JSON findings; edge companion hub with web dashboard"),
        ("Tier 3 — Hospital VPC / Batch", "ROADMAP", TIER3_PURPLE, RGBColor(0xF3, 0xE5, 0xF5),
         "Overnight batch: 27B+ models, rare-disease screens, cohort analytics via FHIR ingest. On-device ScheduledPrognosisWorker is the shipped edge version"),
    ]

    card_y = cy
    for title, status, clr, bg, desc in tiers:
        card_h = 1.55
        # Card background
        add_rect(slide, L_LEFT + PAD, card_y, L_W - PAD * 2, card_h, bg, clr, Pt(2))
        # Status chip
        chip_w = 1.6
        add_rounded_rect(slide, L_LEFT + L_W - PAD - chip_w - 0.1, card_y + 0.1,
                         chip_w, 0.32, clr)
        add_textbox(slide, L_LEFT + L_W - PAD - chip_w - 0.1, card_y + 0.12,
                    chip_w, 0.28, status, font_size=13, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        # Title
        add_textbox(slide, L_LEFT + PAD + 0.15, card_y + 0.1,
                    L_W - PAD * 2 - chip_w - 0.5, 0.35,
                    title, font_size=18, bold=True, color=clr)
        # Desc
        add_textbox(slide, L_LEFT + PAD + 0.15, card_y + 0.5,
                    L_W - PAD * 2 - 0.3, card_h - 0.6,
                    desc, font_size=14, color=BLACK)
        card_y += card_h + 0.15

    # --- Trust Boundary ---
    y += tier_h + GAP
    trust_h = 2.8
    cy = section_block(slide, L_LEFT, y, L_W, trust_h,
                        "Privacy & Trust Boundary", TIER1_GREEN, TRUST_BG, TIER1_GREEN)

    trust_text = (
        "Clinical narrative and chart Q&A default on-device. Edge companion and "
        "cloud tiers are user-governed; chart text does not go to third-party LLMs "
        "unless the clinician explicitly selects Gemini API.\n\n"
        "Public internet is used only for model weight delivery (Hugging Face CDN) "
        "and OAuth sign-in — never for clinical inference on patient data."
    )
    add_textbox(slide, L_LEFT + PAD, cy, L_W - PAD * 2, trust_h - 0.85,
                trust_text, font_size=16, bold=True, color=TIER1_GREEN)

    # ================================================================
    # RIGHT COLUMN
    # ================================================================
    # --- Architecture Diagram ---
    y_r = 5.15
    arch_h = 6.6
    cy = section_block(slide, R_LEFT, y_r, R_W, arch_h,
                        "System Architecture", QC_BLUE, SECTION_BG, BORDER_BLUE)

    arch_img = os.path.join(ASSETS, "architecture_v2.png")
    if os.path.exists(arch_img):
        slide.shapes.add_picture(arch_img,
                                 Inches(R_LEFT + PAD), Inches(cy),
                                 Inches(R_W - PAD * 2), Inches(arch_h - 1.0))

    # --- App Interface Composite ---
    y_r += arch_h + GAP
    app_h = 5.2
    cy = section_block(slide, R_LEFT, y_r, R_W, app_h,
                        "Application Interface", DARK_NAVY, SECTION_BG, DARK_NAVY)

    app_img = os.path.join(ASSETS, "app_composite_v2.png")
    if os.path.exists(app_img):
        slide.shapes.add_picture(app_img,
                                 Inches(R_LEFT + PAD), Inches(cy),
                                 Inches(R_W - PAD * 2), Inches(app_h - 0.85))

    # --- Key Highlights ---
    y_r += app_h + GAP
    hl_h = 7.4
    cy = section_block(slide, R_LEFT, y_r, R_W, hl_h,
                        "Key Highlights & Innovations", QC_BLUE, HIGHLIGHT_BG, BORDER_BLUE)

    highlights = [
        ("On-Device Multimodal AI:", " MedGemma 1.5 4B runs text + vision entirely on Snapdragon via custom llama.cpp JNI. Zero cloud inference for clinical data."),
        ("Quantization Discipline:", " In-repo benchmarks across 5 medical datasets (500 samples each) prove Q8_0 is lossless and Q4_K_M preserves reasoning at 2.4 GB."),
        ("Hybrid Vision Pipeline:", " Heavy imaging offloaded to near-edge GPU; only structured JSON returns. Chart text never leaves device in hybrid mode."),
        ("Edge Companion Hub:", " FastAPI laptop server with patient sync, Ollama/Gemini chat routing, chart processing, and web dashboard."),
        ("Smart Vernacular Injection:", " Zero-overhead regex interceptor translates medical jargon to Telugu/Hindi without degrading reasoning benchmarks."),
        ("Longitudinal Intelligence:", " Room DB stores chronological profiles; AI bridges historical records with current symptoms via chart-grounded prompts."),
        ("FHIR Export:", " On-device LLM-assisted FHIR bundle generation for hospital systems. Encrypted storage + PIN/biometric auth."),
        ("Scheduled Prognosis:", " WorkManager overnight batch generates prognoses for all patients when charging — shipped edge version of hospital batch tier."),
    ]

    hl_y = cy
    for title, desc in highlights:
        txBox = slide.shapes.add_textbox(
            Inches(R_LEFT + PAD + 0.1), Inches(hl_y),
            Inches(R_W - PAD * 2 - 0.2), Inches(0.78))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.space_after = Pt(1)

        r1 = p.add_run()
        _set_run(r1, "▶  ", 14, False, QC_BLUE)
        r2 = p.add_run()
        _set_run(r2, title, 17, True, DARK_NAVY)
        r3 = p.add_run()
        _set_run(r3, desc, 15, False, BLACK)
        hl_y += 0.85

    # --- Benchmarks & Performance ---
    y_r += hl_h + GAP
    kpi_h = 10.5
    cy = section_block(slide, R_LEFT, y_r, R_W, kpi_h,
                        "Benchmarks & Performance", ORANGE, KPI_BG, ORANGE)

    # KPI cards image
    kpi_img = os.path.join(ASSETS, "kpi_cards_v2.png")
    if os.path.exists(kpi_img):
        slide.shapes.add_picture(kpi_img,
                                 Inches(R_LEFT + PAD), Inches(cy),
                                 Inches(R_W - PAD * 2), Inches(1.6))

    # Benchmark chart image
    bench_img = os.path.join(ASSETS, "benchmark_chart_v2.png")
    if os.path.exists(bench_img):
        slide.shapes.add_picture(bench_img,
                                 Inches(R_LEFT + PAD), Inches(cy + 1.9),
                                 Inches(R_W - PAD * 2), Inches(4.2))

    # Trust diagram
    trust_img = os.path.join(ASSETS, "trust_diagram_v2.png")
    if os.path.exists(trust_img):
        slide.shapes.add_picture(trust_img,
                                 Inches(R_LEFT + PAD), Inches(cy + 6.4),
                                 Inches(R_W - PAD * 2), Inches(2.2))

    # Key findings text
    findings_y = cy + 8.8
    findings = [
        "✓  Q8_0 is lossless: zero accuracy drop across all 5 benchmarks at ~47% size reduction",
        "✓  Q6_K is near-lossless: within noise margin on all benchmarks at ~59% size reduction",
        "✓  Q4_K_M trades knowledge recall for zero reasoning drop at 2.4 GB — ideal edge compromise",
    ]
    txBox = slide.shapes.add_textbox(
        Inches(R_LEFT + PAD), Inches(findings_y),
        Inches(R_W - PAD * 2), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.paragraphs[0].text = ""
    for f in findings:
        add_para(tf, f, size=15, bold=False, color=DARK_NAVY, space_after=4)

    # ================================================================
    # FOOTER
    # ================================================================
    add_rect(slide, 0, 34.2, 24, 1.8, DARK_NAVY)
    add_textbox(slide, 0.6, 34.4, 10, 0.5,
                "Itikela Bhaskar  |  Vinay G",
                font_size=26, bold=True, color=WHITE)
    add_textbox(slide, 0.6, 34.9, 10, 0.4,
                "Qualcomm India  ·  SparQ 2026",
                font_size=20, color=ACCENT_GOLD)
    add_textbox(slide, 0.6, 35.3, 10, 0.4,
                "vinayg1752004@gmail.com",
                font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA))

    add_textbox(slide, 14, 34.4, 9.5, 0.5,
                "SparQ 2026 — Pervasive AI Theme",
                font_size=22, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.RIGHT)
    add_textbox(slide, 14, 34.9, 9.5, 0.4,
                "On-device MedGemma 1.5 4B  ·  Edge Companion  ·  FHIR Export  ·  Q4_K_M GGUF",
                font_size=15, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)

    prs.save(OUT_PPTX)
    print(f"Poster v2 saved: {OUT_PPTX}")


if __name__ == "__main__":
    build_poster()
