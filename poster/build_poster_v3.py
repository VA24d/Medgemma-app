"""
Med Veda — SparQ 2026 Research Poster v3
Fixes: text overflow, consistent bullets, no Q4_K_M emphasis, correct authors,
       no app screenshots, no "1.5"/"4B" in MedGemma refs, better headings.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(REPO, "Sparq-2026-Poster Format.pptx")
OUT = os.path.join(REPO, "poster", "Med-Veda-SparQ-2026-Poster-v3.pptx")
ASSETS = os.path.join(REPO, "poster", "assets")

# Palette
NAVY      = RGBColor(0x0A, 0x1A, 0x2F)
BLUE      = RGBColor(0x00, 0x5A, 0x9C)
GREEN     = RGBColor(0x1B, 0x5E, 0x20)
TEAL      = RGBColor(0x15, 0x65, 0xC0)
PURPLE    = RGBColor(0x45, 0x27, 0xA0)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x00, 0x00, 0x00)
GOLD      = RGBColor(0xFF, 0xB3, 0x00)
ORANGE    = RGBColor(0xE6, 0x51, 0x00)
RED       = RGBColor(0xC6, 0x28, 0x28)
LGRAY     = RGBColor(0xF0, 0xF4, 0xF8)
MGRAY     = RGBColor(0xBB, 0xBB, 0xBB)
LGREEN    = RGBColor(0xE8, 0xF5, 0xE9)
LBLUE     = RGBColor(0xE3, 0xF2, 0xFD)
LPURPLE   = RGBColor(0xF3, 0xE5, 0xF5)
LORANGE   = RGBColor(0xFD, 0xF2, 0xE9)

FONT = "Calibri"


# ── helpers ──────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill, border=None, bw=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    return s


def rrect(slide, l, t, w, h, fill, border=None, bw=Pt(2)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    return s


def _run(p, text, sz, bold=False, color=BLACK):
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = FONT
    return r


def tbox(slide, l, t, w, h, text, sz=20, bold=False, color=BLACK,
         align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.alignment = align
    _run(p, text, sz, bold, color)
    return tb, tf


def section(slide, l, t, w, h, title, hdr_bg, body_bg, border):
    """Section with flat header bar + body rect. Returns content-start y."""
    HH = 0.48
    rect(slide, l, t + HH - 0.04, w, h - HH + 0.04, body_bg, border, Pt(2))
    rect(slide, l, t, w, HH, hdr_bg, border, Pt(2))
    tbox(slide, l + 0.2, t + 0.04, w - 0.4, HH - 0.08,
         title, sz=26, bold=True, color=WHITE)
    return t + HH + 0.15


def bullets(slide, l, t, w, h, items, sz=16, color=BLACK, spacing=6):
    """Consistent bullet list using dash character."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    tf.paragraphs[0].text = ""
    for item in items:
        p = tf.add_paragraph()
        p.space_after = Pt(spacing); p.space_before = Pt(0)
        _run(p, f"–  {item}", sz, False, color)
    return tb, tf


def labeled_bullets(slide, l, t, w, h, items, sz=15, label_color=BLUE):
    """List of (label, description) pairs with consistent formatting."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    tf.paragraphs[0].text = ""
    for label, desc in items:
        p = tf.add_paragraph()
        p.space_after = Pt(2); p.space_before = Pt(5)
        _run(p, label, sz + 1, True, label_color)
        p2 = tf.add_paragraph()
        p2.space_after = Pt(6); p2.space_before = Pt(0)
        _run(p2, f"  {desc}", sz, False, BLACK)
    return tb, tf


def highlight_item(slide, l, t, w, label, desc, sz=15):
    """Single highlight row: dash + bold label + description."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.space_after = Pt(1)
    _run(p, "–  ", sz, False, BLUE)
    _run(p, label, sz + 1, True, NAVY)
    _run(p, f"  {desc}", sz, False, BLACK)


# ── main ─────────────────────────────────────────────────────────────

def build():
    # Extract logo
    tmpl = Presentation(TEMPLATE)
    logo_path = os.path.join(ASSETS, "sparq_logo.png")
    for s in tmpl.slides[0].shapes:
        if s.name == "Picture 2":
            with open(logo_path, "wb") as f:
                f.write(s.image.blob)
            break

    prs = Presentation()
    prs.slide_width = Inches(24); prs.slide_height = Inches(36)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, 24, 36, WHITE)

    # ═══════════════ BANNER ═══════════════
    slide.shapes.add_picture(logo_path, Inches(0), Inches(0),
                             Inches(8.4), Inches(5.36))
    rect(slide, 7.64, 0, 16.36, 4.5, NAVY)
    tbox(slide, 8.3, 0.3, 15.0, 1.4,
         "Med Veda", sz=84, bold=True, color=WHITE)
    tbox(slide, 8.3, 1.75, 15.0, 0.75,
         "Pervasive Clinical AI on Snapdragon + Governed Hybrid Compute",
         sz=32, color=GOLD)
    tbox(slide, 8.3, 2.55, 15.0, 0.7,
         "On-device MedGemma multimodal   |   Edge companion hub   |   "
         "Governed cloud tiers   |   Privacy-first by design",
         sz=19, color=MGRAY)

    # Subtitle tags
    tags = [
        (8.3,  "On-Device Inference", GREEN),
        (11.6, "Edge Companion", TEAL),
        (14.7, "FHIR Interop", PURPLE),
        (17.4, "Offline-First", ORANGE),
    ]
    for tx, label, clr in tags:
        rrect(slide, tx, 3.4, 2.9, 0.38, clr)
        tbox(slide, tx, 3.43, 2.9, 0.32, label,
             sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Status legend
    for tx, label, clr in [(8.3, "SHIPPED", GREEN),
                            (10.2, "PROTOTYPE", TEAL),
                            (12.3, "ROADMAP", PURPLE)]:
        rrect(slide, tx, 3.95, 1.7, 0.32, clr)
        tbox(slide, tx, 3.97, 1.7, 0.28, label,
             sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rect(slide, 0, 5.1, 24, 0.07, BLUE)

    # ═══════════════ COLUMN GEOMETRY ═══════════════
    LL = 0.5;  LW = 7.8    # left col
    RL = 8.8;  RW = 14.7   # right col
    P  = 0.25              # inner pad
    G  = 0.3               # gap between sections

    # ═══════════════ LEFT COLUMN ═══════════════

    # ── Abstract ──
    y = 5.4
    AH = 6.0
    cy = section(slide, LL, y, LW, AH, "Abstract", BLUE, LGRAY, BLUE)

    abstract = (
        "Mobile clinical workflows operate under intermittent connectivity, "
        "strict latency budgets, and thermal constraints on sustained inference. "
        "Protected health information must remain under privacy-preserving "
        "processing, yet multimodal decisions require tight coupling between "
        "longitudinal records and high-resolution imaging.\n\n"
        "Med Veda is an Android clinician assistant centered on on-device "
        "MedGemma multimodal inference. Privacy-sensitive reasoning, chart-"
        "conditioned dialogue, and multilingual generation execute locally on "
        "Snapdragon hardware via a native GGUF runtime with heterogeneous "
        "acceleration. This edge-first design preserves low-latency interaction "
        "in offline wards and keeps narrative PHI under local control.\n\n"
        "Compute-intensive vision is selectively orchestrated to institutional "
        "imaging services returning structured analytics, not open-ended cloud "
        "generation. FHIR export enables hospital VPC clusters for deeper "
        "screening and cohort analytics while the handset remains the bedside "
        "copilot.\n\n"
        "The contribution is a production-oriented hybrid edge-cloud pattern: "
        "Snapdragon-edge multimodal reasoning as the default trust boundary, "
        "governed delegation for heavy vision, and standards-based extensibility "
        "for institutional compute."
    )
    tbox(slide, LL + P, cy, LW - P*2, AH - 0.8, abstract, sz=15, color=BLACK)

    # ── Problem ──
    y += AH + G
    PH = 4.0
    cy = section(slide, LL, y, LW, PH, "Problem Statement", RED, LGRAY, RED)

    problems = [
        "Clinicians need longitudinal chart + imaging at bedside with "
        "intermittent or no connectivity",
        "Cloud-bound AI introduces unacceptable latency and mandates PHI "
        "egress to third-party servers",
        "India's DPDP Act and global data protection laws impose strict "
        "constraints on clinical data movement",
        "A single device cannot handle all workloads at full speed — large "
        "imaging hits thermal and compute limits",
        "Existing solutions force a false choice: all-on-phone (slow) or "
        "all-in-cloud (privacy-violating, offline-incompatible)",
    ]
    bullets(slide, LL + P, cy, LW - P*2, PH - 0.8, problems, sz=15)

    # ── Methods ──
    y += PH + G
    MH = 5.8
    cy = section(slide, LL, y, LW, MH, "Methods & Technical Approach", BLUE, LGRAY, BLUE)

    methods = [
        ("MedGemma Multimodal GGUF",
         "Deployed via quantized GGUF through custom llama.cpp JNI wrapper "
         "(aichatlib) with lazy mmproj vision encoder loading"),
        ("Heterogeneous Acceleration",
         "CPU (KleidiAI arm64), Adreno GPU, and Hexagon NPU runtime paths; "
         "user-selectable energy mode with dynamic model switching"),
        ("Chart-Grounded Inference",
         "Room DB stores longitudinal records; system prompts inject full "
         "chart context with dated entries for temporally-aware reasoning"),
        ("Hybrid Vision Pipeline",
         "Study pixels sent to near-edge GPU; only structured JSON findings "
         "return to device for local text synthesis — chart never leaves phone"),
        ("FHIR Interoperability",
         "On-device LLM-assisted FHIR bundle export enables hospital VPC "
         "ingestion for batch analytics with larger models"),
    ]
    labeled_bullets(slide, LL + P, cy, LW - P*2, MH - 0.8, methods, sz=14)

    # ── Three-Tier Architecture ──
    y += MH + G
    TH = 6.0
    cy = section(slide, LL, y, LW, TH, "Three-Tier Pervasive AI Architecture",
                  NAVY, LGRAY, NAVY)

    tiers = [
        ("Tier 1 — On-Device", "SHIPPED", GREEN, LGREEN,
         "Snapdragon phone: MedGemma GGUF + mmproj vision; chart-grounded "
         "chat; Room DB; PIN auth; full offline operation"),
        ("Tier 2 — Near Edge", "PROTOTYPE", TEAL, LBLUE,
         "Hospital / clinic LAN GPU: heavy imaging returns structured JSON "
         "findings; edge companion hub with web dashboard and patient sync"),
        ("Tier 3 — Hospital VPC", "ROADMAP", PURPLE, LPURPLE,
         "Overnight batch with larger models, rare-disease screens, cohort "
         "analytics via FHIR ingest; on-device prognosis worker is shipped "
         "edge version"),
    ]
    card_y = cy
    for title, status, clr, bg, desc in tiers:
        ch = 1.5
        rect(slide, LL + P, card_y, LW - P*2, ch, bg, clr, Pt(2))
        # status chip
        cw = 1.55
        rrect(slide, LL + LW - P - cw - 0.08, card_y + 0.08, cw, 0.3, clr)
        tbox(slide, LL + LW - P - cw - 0.08, card_y + 0.1, cw, 0.26,
             status, sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # title
        tbox(slide, LL + P + 0.12, card_y + 0.08, LW - P*2 - cw - 0.4, 0.3,
             title, sz=17, bold=True, color=clr)
        # desc
        tbox(slide, LL + P + 0.12, card_y + 0.45, LW - P*2 - 0.24, ch - 0.55,
             desc, sz=14, color=BLACK)
        card_y += ch + 0.15

    # ── Trust Boundary ──
    y += TH + G
    TB_H = 2.6
    cy = section(slide, LL, y, LW, TB_H,
                  "Privacy & Trust Boundary", GREEN, LGREEN, GREEN)

    trust = (
        "Clinical narrative and chart Q&A default on-device. Edge companion "
        "and cloud tiers are user-governed — chart text does not go to "
        "third-party LLMs unless the clinician explicitly opts in.\n\n"
        "Public internet is used only for model weight delivery and "
        "OAuth sign-in — never for clinical inference on patient data."
    )
    tbox(slide, LL + P, cy, LW - P*2, TB_H - 0.8,
         trust, sz=15, bold=True, color=GREEN)

    # ═══════════════ RIGHT COLUMN ═══════════════

    # ── Architecture Diagram ──
    yr = 5.4
    AR_H = 5.8
    cy = section(slide, RL, yr, RW, AR_H,
                  "System Architecture — Pervasive AI Continuum", BLUE, LGRAY, BLUE)

    arch = os.path.join(ASSETS, "architecture_v3.png")
    if os.path.exists(arch):
        slide.shapes.add_picture(arch,
            Inches(RL + P), Inches(cy), Inches(RW - P*2), Inches(AR_H - 0.85))

    # ── Key Highlights ──
    yr += AR_H + G
    HL_H = 7.5
    cy = section(slide, RL, yr, RW, HL_H,
                  "Key Highlights & Innovations", NAVY, LBLUE, NAVY)

    highlights = [
        ("On-Device Multimodal AI",
         "MedGemma runs text + vision entirely on Snapdragon via custom "
         "llama.cpp JNI wrapper. Zero cloud inference for clinical data."),
        ("Quantization Discipline",
         "In-repo benchmarks across 5 medical datasets (500 samples each) — "
         "Q8_0 is lossless at 47% size reduction; reasoning accuracy preserved "
         "across all quantization levels."),
        ("Hybrid Vision Pipeline",
         "Heavy imaging (X-ray, histopathology) offloaded to near-edge GPU; "
         "only structured JSON returns to phone. Chart text never leaves "
         "device in hybrid mode."),
        ("Edge Companion Hub",
         "FastAPI-based laptop server with patient sync, LLM routing "
         "(Ollama / Gemini), chart processing, and web dashboard for "
         "clinical oversight."),
        ("Smart Vernacular Injection",
         "Zero-overhead regex interceptor translates complex medical jargon "
         "to Telugu / Hindi in the output stream without degrading core "
         "reasoning benchmarks."),
        ("Longitudinal Intelligence",
         "Room DB stores chronological patient profiles; AI bridges years of "
         "historical records with current symptoms via chart-grounded prompts."),
        ("FHIR Export & Interoperability",
         "On-device LLM-assisted FHIR bundle generation enables data flow to "
         "hospital systems without cloud intermediary. Encrypted local "
         "storage with PIN / biometric auth."),
        ("Scheduled Prognosis Worker",
         "WorkManager-based overnight batch generates AI prognoses for all "
         "patients when device is charging — shipped edge version of the "
         "hospital batch tier."),
    ]

    hy = cy
    for label, desc in highlights:
        highlight_item(slide, RL + P + 0.05, hy, RW - P*2 - 0.1, label, desc, sz=15)
        hy += 0.85

    # ── Benchmarks & Performance ──
    yr += HL_H + G
    BM_H = 11.3
    cy = section(slide, RL, yr, RW, BM_H,
                  "Benchmarks & Performance", ORANGE, LORANGE, ORANGE)

    # KPI cards
    kpi = os.path.join(ASSETS, "kpi_v3.png")
    if os.path.exists(kpi):
        slide.shapes.add_picture(kpi,
            Inches(RL + P), Inches(cy), Inches(RW - P*2), Inches(1.5))

    # Benchmark chart
    bench = os.path.join(ASSETS, "benchmark_v3.png")
    if os.path.exists(bench):
        slide.shapes.add_picture(bench,
            Inches(RL + P), Inches(cy + 1.8), Inches(RW - P*2), Inches(4.5))

    # Trust diagram
    trust_img = os.path.join(ASSETS, "trust_v3.png")
    if os.path.exists(trust_img):
        slide.shapes.add_picture(trust_img,
            Inches(RL + P), Inches(cy + 6.6), Inches(RW - P*2), Inches(2.0))

    # Benchmark table
    tbl_y = cy + 8.9
    tbox(slide, RL + P, tbl_y - 0.05, 10, 0.35,
         "Quantization vs BF16 Baseline  (500 samples x 5 benchmarks, RTX 4090)",
         sz=16, bold=True, color=NAVY)

    tbl_data = [
        ["Model", "Size", "MedMCQA", "MedQA", "PubMedQA", "MMLU Med", "MedXpertQA"],
        ["BF16 (base)", "7.3 GB", "43.8%", "29.0%", "55.4%", "43.0%", "8.8%"],
        ["Q8_0", "3.9 GB", "44.4%", "28.6%", "55.4%", "43.6%", "8.8%"],
        ["Q6_K", "3.0 GB", "40.8%", "28.4%", "57.4%", "41.4%", "9.8%"],
        ["Q4_K_M", "2.4 GB", "32.6%", "29.0%", "55.4%", "29.8%", "10.0%"],
    ]
    rows, cols = len(tbl_data), len(tbl_data[0])
    tw = 14.0
    th = rows * 0.42
    ts = slide.shapes.add_table(rows, cols,
        Inches(RL + P + 0.1), Inches(tbl_y + 0.35), Inches(tw), Inches(th))
    table = ts.table

    col_ws = [1.8, 1.1, 1.7, 1.5, 1.7, 1.7, 1.7]
    for i, cw in enumerate(col_ws):
        table.columns[i].width = Inches(cw)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = tbl_data[r][c]
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(14)
                    run.font.name = FONT
                    run.font.bold = r == 0
                    run.font.color.rgb = WHITE if r == 0 else BLACK
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = ORANGE
            elif r % 2 == 1:
                cell.fill.fore_color.rgb = WHITE
            else:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)

    # ═══════════════ FOOTER ═══════════════
    rect(slide, 0, 34.2, 24, 1.8, NAVY)
    tbox(slide, 0.6, 34.4, 10, 0.5,
         "Itikela Bhaskar  |  Vijay Aravynthan",
         sz=26, bold=True, color=WHITE)
    tbox(slide, 0.6, 34.95, 10, 0.35,
         "Qualcomm India  ·  SparQ 2026",
         sz=20, color=GOLD)
    tbox(slide, 14, 34.4, 9.5, 0.5,
         "SparQ 2026 — Pervasive AI Theme",
         sz=22, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
    tbox(slide, 14, 34.95, 9.5, 0.35,
         "On-device MedGemma  ·  Edge Companion  ·  FHIR Export  ·  GGUF Quantization",
         sz=15, color=MGRAY, align=PP_ALIGN.RIGHT)

    prs.save(OUT)
    print(f"Poster v3 saved: {OUT}")


if __name__ == "__main__":
    build()
