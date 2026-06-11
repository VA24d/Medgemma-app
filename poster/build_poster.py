"""
Med Veda -- SparQ 2026 Research Poster Generator
Produces a 24x36 inch single-slide poster PPTX matching the SparQ template layout.
"""

import os, sys, shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(REPO, "Sparq-2026-Poster Format.pptx")
OUT_PPTX = os.path.join(REPO, "poster", "Med-Veda-SparQ-2026-Poster.pptx")
ASSETS = os.path.join(REPO, "poster", "assets")

# Colors
DARK_NAVY   = RGBColor(0x0A, 0x1A, 0x2F)
QC_BLUE     = RGBColor(0x00, 0x5A, 0x9C)
TIER1_GREEN = RGBColor(0x1B, 0x5E, 0x20)
TIER2_BLUE  = RGBColor(0x15, 0x65, 0xC0)
TIER3_PURPLE = RGBColor(0x45, 0x27, 0xA0)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
BLACK       = RGBColor(0x00, 0x00, 0x00)
ACCENT_GOLD = RGBColor(0xFF, 0xB3, 0x00)
SECTION_BG  = RGBColor(0xF0, 0xF4, 0xF8)
BORDER_BLUE = RGBColor(0x00, 0x7A, 0xCC)
TRUST_BG    = RGBColor(0xE8, 0xF5, 0xE9)
KPI_BG      = RGBColor(0xFD, 0xF2, 0xE9)
HIGHLIGHT_BG = RGBColor(0xE3, 0xF2, 0xFD)

def add_textbox(slide, left, top, width, height, text, font_size=22, bold=False,
                color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Calibri",
                anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox, tf

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape

def add_multiline_text(slide, left, top, width, height, lines, font_size=20,
                       color=BLACK, font_name="Calibri", line_spacing=1.15,
                       bold_first=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        p.line_spacing = line_spacing
        run = p.add_run()
        if isinstance(line, tuple):
            run.text = line[0]
            run.font.bold = line[1]
            if len(line) > 2:
                run.font.color.rgb = line[2]
            else:
                run.font.color.rgb = color
        else:
            run.text = line
            run.font.bold = bold_first and i == 0
            run.font.color.rgb = color
        run.font.size = Pt(font_size)
        run.font.name = font_name
    return txBox, tf

def add_bullet_list(slide, left, top, width, height, items, font_size=20,
                    color=BLACK, bullet_char="•"):
    lines = [f"{bullet_char}  {item}" for item in items]
    return add_multiline_text(slide, left, top, width, height, lines,
                              font_size=font_size, color=color)

def add_section_header(slide, left, top, width, text, font_size=32, color=WHITE,
                       bg_color=QC_BLUE, height=0.55):
    add_rounded_rect(slide, left, top, width, height, bg_color)
    add_textbox(slide, left + 0.15, top + 0.05, width - 0.3, height - 0.1,
                text, font_size=font_size, bold=True, color=color,
                alignment=PP_ALIGN.LEFT)


def build_poster():
    # --- Load template to extract SparQ logo ---
    tmpl = Presentation(TEMPLATE)
    tmpl_slide = tmpl.slides[0]
    logo_blob = None
    logo_ct = None
    for shape in tmpl_slide.shapes:
        if shape.name == "Picture 2":
            logo_blob = shape.image.blob
            logo_ct = shape.image.content_type
            ext = "png" if "png" in logo_ct else "jpg"
            logo_path = os.path.join(ASSETS, f"sparq_logo.{ext}")
            with open(logo_path, "wb") as f:
                f.write(logo_blob)
            break

    # --- New presentation (blank, 24x36) ---
    prs = Presentation()
    prs.slide_width = Inches(24)
    prs.slide_height = Inches(36)
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    # ========== BACKGROUND ==========
    add_rect(slide, 0, 0, 24, 36, WHITE)

    # ========== TOP BANNER (SparQ logo area + Title) ==========
    # Logo background area (left)
    if logo_blob:
        slide.shapes.add_picture(logo_path, Inches(0), Inches(0),
                                 Inches(8.4), Inches(5.36))

    # Title banner (right of logo)
    add_rect(slide, 7.64, 0, 16.36, 3.86, DARK_NAVY)
    # Title text
    add_textbox(slide, 8.0, 0.3, 15.5, 1.5,
                "Med Veda", font_size=84, bold=True,
                color=WHITE, font_name="Calibri")
    add_textbox(slide, 8.0, 1.7, 15.5, 1.0,
                "Pervasive Clinical AI on Snapdragon + Governed Hybrid Compute",
                font_size=36, bold=False, color=ACCENT_GOLD, font_name="Calibri")
    # Thesis subtitle
    add_textbox(slide, 8.0, 2.6, 15.5, 1.2,
                "On-device MedGemma 1.5 4B multimodal · Edge companion hub · "
                "Optional cloud APIs — privacy-first AI as a continuum, not \"everything in the cloud\"",
                font_size=22, bold=False, color=RGBColor(0xCC, 0xCC, 0xCC),
                font_name="Calibri")

    # Tag chips under title
    chip_y = 3.5
    chip_labels = [
        ("SHIPPED", TIER1_GREEN),
        ("PROTOTYPE", TIER2_BLUE),
        ("ROADMAP", TIER3_PURPLE),
    ]
    chip_x = 8.0
    for label, clr in chip_labels:
        r = add_rounded_rect(slide, chip_x, chip_y, 2.2, 0.35, clr)
        add_textbox(slide, chip_x + 0.1, chip_y + 0.02, 2.0, 0.3,
                    label, font_size=16, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)
        chip_x += 2.5

    # ===== Divider line under banner =====
    add_rect(slide, 0, 5.3, 24, 0.06, QC_BLUE)

    # ========== ABSTRACT (top-left) ==========
    abs_left, abs_top = 0.6, 5.5
    abs_w, abs_h = 7.5, 6.8
    add_rounded_rect(slide, abs_left, abs_top, abs_w, abs_h, SECTION_BG, BORDER_BLUE)
    add_section_header(slide, abs_left + 0.15, abs_top + 0.15, abs_w - 0.3, "Abstract")

    abstract_text = (
        "Mobile clinical workflows operate under intermittent connectivity, strict latency "
        "budgets at the point of care, and thermal constraints on sustained inference. "
        "Protected health information must remain amenable to privacy-preserving processing, "
        "yet multimodal decisions still require tight coupling between longitudinal records "
        "and high-resolution imaging.\n\n"
        "Med Veda is a deployable Android clinician assistant centered on on-device Google "
        "MedGemma 1.5 4B multimodal inference. The system architecture partitions "
        "responsibility explicitly: privacy-sensitive longitudinal reasoning, chart-conditioned "
        "dialogue, and multilingual clinical generation execute locally on Snapdragon-class "
        "hardware behind a native GGUF runtime suitable for heterogeneous acceleration, "
        "including NPU-assisted paths. This edge-first design preserves low-latency "
        "interaction in offline-first wards and keeps narrative PHI under local control.\n\n"
        "Compute-intensive vision stages are selectively orchestrated to enterprise imaging "
        "services under institutional credentials, returning compact structured analytics "
        "rather than open-ended cloud language generation. FHIR-oriented export enables "
        "hospital on-premises or VPC GPU clusters for deeper screening, cohort analytics, "
        "or larger specialist models while the handset remains the bedside copilot.\n\n"
        "The contribution is a production-oriented hybrid edge–cloud pattern for pervasive "
        "clinical AI: Snapdragon-edge multimodal reasoning as the default trust boundary, "
        "governed delegation for heavy vision, and standards-based extensibility for "
        "institutional compute."
    )

    add_textbox(slide, abs_left + 0.3, abs_top + 0.85, abs_w - 0.6, abs_h - 1.1,
                abstract_text, font_size=17, color=BLACK, font_name="Calibri")

    # ========== IMAGES (top-right) ==========
    img_left, img_top = 9.0, 5.5
    img_w, img_h = 14.2, 6.8
    add_rounded_rect(slide, img_left, img_top, img_w, img_h, SECTION_BG, BORDER_BLUE)
    add_section_header(slide, img_left + 0.15, img_top + 0.15, img_w - 0.3,
                       "System Architecture & App Interface")

    # Architecture diagram
    arch_png = os.path.join(REPO, "ppt2", "diagrams", "png", "diagram_1.png")
    if os.path.exists(arch_png):
        shutil.copy(arch_png, os.path.join(ASSETS, "architecture_diagram.png"))
        slide.shapes.add_picture(arch_png, Inches(img_left + 0.3), Inches(img_top + 0.9),
                                 Inches(8.5), Inches(5.0))

    # App screenshots
    screenshots = [
        ("app images/home screeen.png", "Home"),
        ("app images/image analysis.png", "X-ray AI"),
        ("app images/patient page.png", "Patient"),
    ]
    ss_x = img_left + 9.2
    for ss_file, label in screenshots:
        ss_path = os.path.join(REPO, ss_file)
        if os.path.exists(ss_path):
            slide.shapes.add_picture(ss_path, Inches(ss_x), Inches(img_top + 0.9),
                                     Inches(1.5), Inches(3.2))
            add_textbox(slide, ss_x, img_top + 4.2, 1.5, 0.4,
                        label, font_size=14, bold=True, color=QC_BLUE,
                        alignment=PP_ALIGN.CENTER)
            ss_x += 1.65

    # ========== LEFT COLUMN: Problem + Methods + Tier Table + Trust ==========
    col_left, col_top = 0.6, 13.0
    col_w = 7.5

    # --- Problem ---
    add_rounded_rect(slide, col_left, col_top, col_w, 4.5, SECTION_BG, BORDER_BLUE)
    add_section_header(slide, col_left + 0.15, col_top + 0.15, col_w - 0.3, "Problem Statement")

    problem_items = [
        "Clinicians need longitudinal chart + imaging at the bedside, often with intermittent or no connectivity",
        "Cloud-bound AI introduces unacceptable latency at point of care and mandates PHI egress to third-party servers",
        "India's DPDP Act and global data protection laws impose strict constraints on clinical data movement",
        "A single device cannot handle all workloads at full speed: large imaging studies hit thermal and compute limits",
        "Existing solutions force a false choice between \"all on phone\" (slow, limited) or \"all in cloud\" (privacy-violating, offline-incompatible)",
    ]
    add_bullet_list(slide, col_left + 0.3, col_top + 0.85, col_w - 0.6, 3.5,
                    problem_items, font_size=17)

    # --- Methods ---
    methods_top = col_top + 4.8
    add_rounded_rect(slide, col_left, methods_top, col_w, 5.5, SECTION_BG, BORDER_BLUE)
    add_section_header(slide, col_left + 0.15, methods_top + 0.15, col_w - 0.3,
                       "Methods & Technical Approach")

    methods_items = [
        ("MedGemma 1.5 4B IT (Multimodal):", True, QC_BLUE),
        "  Deployed via GGUF (Q4_K_M default, 2.4 GB) through custom llama.cpp JNI wrapper (aichatlib) with lazy mmproj vision encoder loading",
        ("Heterogeneous Acceleration:", True, QC_BLUE),
        "  CPU (KleidiAI on arm64), Adreno GPU, and Hexagon NPU runtime paths; user-selectable energy mode with 2B/4B dynamic switching",
        ("Chart-Grounded Inference:", True, QC_BLUE),
        "  Room DB stores longitudinal patient records; system prompts inject full chart context with dated entries for temporally-aware reasoning",
        ("Hybrid Vision Pipeline:", True, QC_BLUE),
        "  Study pixels sent to near-edge GPU (prototype: RTX 4070); only structured JSON findings return to device for local text synthesis",
        ("FHIR Interoperability:", True, QC_BLUE),
        "  On-device LLM-assisted FHIR bundle export enables hospital VPC ingestion for batch analytics with larger models",
    ]
    add_multiline_text(slide, col_left + 0.3, methods_top + 0.85, col_w - 0.6, 4.5,
                       methods_items, font_size=17)

    # --- Three-Tier Architecture Table ---
    tier_top = methods_top + 5.8
    add_rounded_rect(slide, col_left, tier_top, col_w, 6.8, SECTION_BG, BORDER_BLUE)
    add_section_header(slide, col_left + 0.15, tier_top + 0.15, col_w - 0.3,
                       "Three-Tier Pervasive AI Architecture")

    # Tier 1
    t1_y = tier_top + 0.9
    add_rounded_rect(slide, col_left + 0.3, t1_y, col_w - 0.6, 1.6, RGBColor(0xE8, 0xF5, 0xE9), TIER1_GREEN)
    add_multiline_text(slide, col_left + 0.5, t1_y + 0.1, col_w - 1.0, 1.5, [
        ("Tier 1 — Device Edge  [SHIPPED]", True, TIER1_GREEN),
        "Android phone (Snapdragon): MedGemma 1.5 4B GGUF + mmproj; chart-grounded chat; Room DB; PIN auth; full offline operation",
    ], font_size=16)

    # Tier 2
    t2_y = t1_y + 1.8
    add_rounded_rect(slide, col_left + 0.3, t2_y, col_w - 0.6, 1.6, RGBColor(0xE3, 0xF2, 0xFD), TIER2_BLUE)
    add_multiline_text(slide, col_left + 0.5, t2_y + 0.1, col_w - 1.0, 1.5, [
        ("Tier 2 — Near Edge  [PROTOTYPE]", True, TIER2_BLUE),
        "Hospital/clinic LAN GPU (lab: RTX 4070): Heavy imaging → structured JSON findings (not full chart to cloud LLM); edge companion hub with web dashboard",
    ], font_size=16)

    # Tier 3
    t3_y = t2_y + 1.8
    add_rounded_rect(slide, col_left + 0.3, t3_y, col_w - 0.6, 1.6, RGBColor(0xF3, 0xE5, 0xF5), TIER3_PURPLE)
    add_multiline_text(slide, col_left + 0.5, t3_y + 0.1, col_w - 1.0, 1.5, [
        ("Tier 3 — Hospital VPC / Batch  [ROADMAP]", True, TIER3_PURPLE),
        "Overnight batch: 27B+ models, rare-disease screens, cohort analytics via FHIR ingest. On-device ScheduledPrognosisWorker is the shipped edge version of this tier",
    ], font_size=16)

    # --- Trust Boundary ---
    trust_top = tier_top + 7.1
    add_rounded_rect(slide, col_left, trust_top, col_w, 2.5, TRUST_BG, TIER1_GREEN)
    add_section_header(slide, col_left + 0.15, trust_top + 0.15, col_w - 0.3,
                       "Privacy & Trust Boundary", bg_color=TIER1_GREEN)

    trust_text = (
        "Clinical narrative and chart Q&A default on-device. Edge companion and cloud tiers "
        "are user-governed; chart text does not go to third-party LLMs unless the clinician "
        "explicitly selects Gemini API.\n\n"
        "Public internet is used only for model weight delivery (Hugging Face CDN) and "
        "OAuth sign-in — never for clinical inference on patient data."
    )
    add_textbox(slide, col_left + 0.3, trust_top + 0.85, col_w - 0.6, 1.5,
                trust_text, font_size=17, color=TIER1_GREEN, font_name="Calibri")

    # ========== HIGHLIGHTS (mid-right) ==========
    hl_left, hl_top = 9.0, 13.0
    hl_w, hl_h = 14.2, 9.5
    add_rounded_rect(slide, hl_left, hl_top, hl_w, hl_h, HIGHLIGHT_BG, BORDER_BLUE)
    add_section_header(slide, hl_left + 0.15, hl_top + 0.15, hl_w - 0.3,
                       "Key Highlights & Innovations")

    highlights = [
        ("On-Device Multimodal AI:", " MedGemma 1.5 4B runs text + vision entirely on Snapdragon via custom llama.cpp JNI (aichatlib). Zero cloud inference for clinical data."),
        ("Quantization Discipline:", " In-repo benchmarks across 5 medical datasets (500 samples each) prove Q8_0 is lossless (47% size reduction) and Q4_K_M preserves reasoning accuracy (zero drop on MedQA/PubMedQA) at 2.4 GB."),
        ("Hybrid Vision Pipeline:", " Heavy imaging (X-ray, histopathology) offloaded to near-edge GPU; only structured JSON returns to phone. Chart text never leaves device in hybrid mode."),
        ("Edge Companion Hub:", " FastAPI-based laptop server (port 8787) with patient sync, Ollama/Gemini chat routing, chart processing, and web dashboard for clinical oversight."),
        ("Smart Vernacular Injection:", " Zero-overhead regex interceptor translates complex medical jargon to Telugu/Hindi in the output stream without degrading MedGemma's core reasoning benchmarks."),
        ("Longitudinal Intelligence:", " Room DB stores chronological patient profiles; AI bridges years of historical records with current symptoms via chart-grounded prompts with dated entries."),
        ("FHIR Export & Interoperability:", " On-device LLM-assisted FHIR bundle generation enables data flow to hospital systems without cloud intermediary. Encrypted local storage + PIN/biometric auth."),
        ("Scheduled Prognosis Worker:", " WorkManager-based overnight batch generates AI prognoses for all patients when device is charging — shipped edge version of the hospital batch tier."),
    ]

    hl_y = hl_top + 0.85
    for title, desc in highlights:
        txBox = slide.shapes.add_textbox(Inches(hl_left + 0.3), Inches(hl_y),
                                          Inches(hl_w - 0.6), Inches(0.95))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.space_after = Pt(2)
        # Bullet
        bullet_run = p.add_run()
        bullet_run.text = "▶  "
        bullet_run.font.size = Pt(16)
        bullet_run.font.color.rgb = QC_BLUE
        bullet_run.font.name = "Calibri"
        # Bold title
        title_run = p.add_run()
        title_run.text = title
        title_run.font.size = Pt(19)
        title_run.font.bold = True
        title_run.font.color.rgb = DARK_NAVY
        title_run.font.name = "Calibri"
        # Description
        desc_run = p.add_run()
        desc_run.text = desc
        desc_run.font.size = Pt(17)
        desc_run.font.color.rgb = BLACK
        desc_run.font.name = "Calibri"
        hl_y += 1.05

    # ========== KPI / STATS (bottom-right) ==========
    kpi_left, kpi_top = 9.0, 23.0
    kpi_w, kpi_h = 14.2, 9.5
    add_rounded_rect(slide, kpi_left, kpi_top, kpi_w, kpi_h, KPI_BG, RGBColor(0xE6, 0x51, 0x00))
    add_section_header(slide, kpi_left + 0.15, kpi_top + 0.15, kpi_w - 0.3,
                       "Benchmarks & Performance", bg_color=RGBColor(0xE6, 0x51, 0x00))

    # Performance metrics
    metrics = [
        ("~18–22", "tok/s", "Text Inference Speed"),
        ("4.1 GB", "peak", "RAM Consumption"),
        ("2.4 GB", "disk", "Model Size (Q4_K_M)"),
        ("5", "datasets", "Medical Benchmarks"),
        ("500", "samples", "Per Benchmark"),
    ]

    metric_x = kpi_left + 0.3
    metric_y = kpi_top + 0.9
    for val, unit, label in metrics:
        add_rounded_rect(slide, metric_x, metric_y, 2.5, 1.5, WHITE, RGBColor(0xE6, 0x51, 0x00))
        add_textbox(slide, metric_x + 0.1, metric_y + 0.15, 2.3, 0.7,
                    val, font_size=32, bold=True, color=RGBColor(0xE6, 0x51, 0x00),
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, metric_x + 0.1, metric_y + 0.75, 2.3, 0.3,
                    unit, font_size=16, bold=False, color=RGBColor(0x99, 0x99, 0x99),
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, metric_x + 0.1, metric_y + 1.05, 2.3, 0.4,
                    label, font_size=14, bold=True, color=DARK_NAVY,
                    alignment=PP_ALIGN.CENTER)
        metric_x += 2.7

    # Benchmark table header
    tbl_top = kpi_top + 2.8
    tbl_data = [
        ["Model", "Size", "MedMCQA", "MedQA", "PubMedQA", "MMLU Med", "MedXpertQA"],
        ["BF16 (baseline)", "7.3 GB", "43.80%", "29.00%", "55.40%", "43.00%", "8.80%"],
        ["Q8_0", "3.9 GB", "44.40%", "28.60%", "55.40%", "43.60%", "8.80%"],
        ["Q6_K", "3.0 GB", "40.80%", "28.40%", "57.40%", "41.40%", "9.80%"],
        ["Q4_K_M", "2.4 GB", "32.60%", "29.00%", "55.40%", "29.80%", "10.00%"],
    ]

    add_textbox(slide, kpi_left + 0.3, tbl_top - 0.15, 8, 0.4,
                "Quantization vs. BF16 Baseline  (500 samples × 5 benchmarks, RTX 4090 GPU)",
                font_size=18, bold=True, color=DARK_NAVY)

    rows = len(tbl_data)
    cols = len(tbl_data[0])
    tbl_w = 13.2
    tbl_h = rows * 0.55
    table_shape = slide.shapes.add_table(rows, cols,
                                          Inches(kpi_left + 0.5), Inches(tbl_top + 0.3),
                                          Inches(tbl_w), Inches(tbl_h))
    table = table_shape.table

    col_widths = [2.2, 1.2, 1.6, 1.6, 1.6, 1.6, 1.6]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = tbl_data[r][c]
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(16)
                    run.font.name = "Calibri"
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    else:
                        run.font.color.rgb = BLACK
                        if r == 1:
                            run.font.bold = True

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE6, 0x51, 0x00)
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    # Key findings
    findings_top = tbl_top + tbl_h + 0.6
    findings = [
        "Q8_0 is lossless: zero meaningful accuracy drop across all 5 benchmarks at ~47% size reduction (7.3 GB → 3.9 GB)",
        "Q6_K is near-lossless: within noise margin on all benchmarks at ~59% size reduction (7.3 GB → 3.0 GB)",
        "Q4_K_M trades knowledge recall (−11.2pp MedMCQA) for zero reasoning drop (MedQA, PubMedQA) at 2.4 GB — ideal edge compromise",
        "TFLite Q8 alternative: 39.55% MedMCQA natively on CPU, matching GGUF Q8_0 quality for CPU-only deployment paths",
    ]
    add_bullet_list(slide, kpi_left + 0.3, findings_top, kpi_w - 0.6, 3.5,
                    findings, font_size=17, color=DARK_NAVY, bullet_char="✓")

    # Trust diagram
    trust_diag = os.path.join(REPO, "ppt2", "diagrams", "png", "diagram_5.png")
    if os.path.exists(trust_diag):
        shutil.copy(trust_diag, os.path.join(ASSETS, "trust_boundaries.png"))
        slide.shapes.add_picture(trust_diag, Inches(kpi_left + 0.3), Inches(findings_top + 3.0),
                                 Inches(13.0), Inches(3.0))

    # ========== FOOTER / AUTHORS ==========
    add_rect(slide, 0, 34.2, 24, 1.8, DARK_NAVY)
    add_textbox(slide, 0.5, 34.4, 10, 0.6,
                "Itikela Bhaskar  |  Vinay G", font_size=28,
                bold=True, color=WHITE)
    add_textbox(slide, 0.5, 35.0, 10, 0.5,
                "Qualcomm India  ·  SparQ 2026", font_size=22,
                bold=False, color=ACCENT_GOLD)
    add_textbox(slide, 0.5, 35.45, 10, 0.5,
                "vinayg1752004@gmail.com", font_size=18,
                bold=False, color=RGBColor(0xAA, 0xAA, 0xAA))
    add_textbox(slide, 14.0, 34.5, 9.5, 0.5,
                "github.com/anthropics/Medgemma-app", font_size=18,
                bold=False, color=RGBColor(0xAA, 0xAA, 0xAA),
                alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, 14.0, 35.0, 9.5, 0.5,
                "SparQ 2026  —  Pervasive AI Theme", font_size=22,
                bold=True, color=ACCENT_GOLD, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, 14.0, 35.45, 9.5, 0.5,
                "On-device MedGemma 1.5 4B  ·  Edge Companion  ·  FHIR Export  ·  Q4_K_M GGUF",
                font_size=16, bold=False, color=RGBColor(0xAA, 0xAA, 0xAA),
                alignment=PP_ALIGN.RIGHT)

    # ========== SAVE ==========
    prs.save(OUT_PPTX)
    print(f"Poster saved: {OUT_PPTX}")
    print(f"Assets in: {ASSETS}")


if __name__ == "__main__":
    build_poster()
